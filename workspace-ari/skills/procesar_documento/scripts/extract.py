"""Extraccion de contenido por formato.

Devuelve siempre la tupla (text_blocks, tables, fmt):
  text_blocks: list[{"title": str, "text": str}]
  tables:      list[{"title": str, "rows": list[list]}]
  fmt:         extension del formato de origen

Las librerias pesadas se importan de forma perezosa: si falta una, solo falla el
formato correspondiente.
"""
from __future__ import annotations

import csv
from pathlib import Path

from clean import remove_repeated_headers_footers


def extract(path: Path):
    ext = path.suffix.lower().lstrip(".")
    if ext in ("txt", "md"):
        return _textlike(path), [], ext
    if ext == "csv":
        return [], _csv(path), ext
    if ext in ("xlsx", "xls"):
        return [], _xlsx(path), ext
    if ext == "pdf":
        return _pdf(path)
    if ext == "docx":
        return _docx(path)
    if ext == "pptx":
        return _pptx(path)
    if ext in ("png", "jpg", "jpeg", "tiff", "bmp"):
        return _image(path), [], ext
    raise ValueError(f"Formato no soportado: .{ext}")


def _textlike(path):
    return [{"title": path.stem, "text": path.read_text(encoding="utf-8", errors="replace")}]


def _csv(path):
    with open(path, newline="", encoding="utf-8", errors="replace") as fh:
        sample = fh.read(4096)
        fh.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample)
        except csv.Error:
            dialect = csv.excel
        rows = list(csv.reader(fh, dialect))
    return [{"title": path.stem, "rows": rows}]


def _xlsx(path):
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("Falta 'openpyxl' para leer .xlsx (pip install openpyxl)")
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    tables = []
    for ws in wb.worksheets:
        rows = [["" if c is None else c for c in row]
                for row in ws.iter_rows(values_only=True)
                if any(c is not None for c in row)]
        if rows:
            tables.append({"title": ws.title, "rows": rows})
    wb.close()
    return tables


def _pdf(path):
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("Falta 'pdfplumber' para leer .pdf (pip install pdfplumber)")
    pages_text, tables = [], []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            txt = page.extract_text() or ""
            if txt.strip():
                pages_text.append(txt)
            for j, tbl in enumerate(page.extract_tables() or [], 1):
                rows = [["" if c is None else c for c in r] for r in tbl]
                if rows:
                    tables.append({"title": f"pagina{i}-tabla{j}", "rows": rows})
    pages_text = remove_repeated_headers_footers(pages_text)
    body = "\n\n".join(pages_text)
    return ([{"title": "cuerpo", "text": body}] if body.strip() else []), tables, "pdf"


def _docx(path):
    try:
        import docx  # python-docx
    except ImportError:
        raise RuntimeError("Falta 'python-docx' para leer .docx (pip install python-docx)")
    document = docx.Document(str(path))
    parts = []
    for p in document.paragraphs:
        style = (p.style.name or "").lower() if p.style else ""
        text = p.text.strip()
        if not text:
            continue
        if style.startswith("heading 1") or style.startswith("titulo 1"):
            parts.append(f"# {text}")
        elif style.startswith("heading") or style.startswith("titulo"):
            parts.append(f"## {text}")
        else:
            parts.append(text)
    tables = []
    for k, t in enumerate(document.tables, 1):
        rows = [[cell.text for cell in row.cells] for row in t.rows]
        if rows:
            tables.append({"title": f"tabla{k}", "rows": rows})
    body = "\n\n".join(parts)
    return ([{"title": "cuerpo", "text": body}] if body.strip() else []), tables, "docx"


def _pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("Falta 'python-pptx' para leer .pptx (pip install python-pptx)")
    prs = Presentation(str(path))
    blocks, tables = [], []
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[c.text for c in row.cells] for row in shape.table.rows]
                if rows:
                    tables.append({"title": f"diapositiva{i}-tabla", "rows": rows})
            elif shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    lines.append(t)
        if lines:
            blocks.append({"title": f"Diapositiva {i}", "text": "\n".join(lines)})
    return blocks, tables, "pptx"


def _image(path):
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        raise RuntimeError("Falta 'pytesseract'+'Pillow' para OCR (pip install pytesseract pillow)")
    return [{"title": "ocr", "text": pytesseract.image_to_string(Image.open(path), lang="spa+eng")}]